# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
from openai import OpenAI
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import os

# ===================== 【你的真实配置 100% 精准填写】 =====================
API_KEY = "ark-4aaa9335-e5ba-4f88-ae37-b0a0ca396ca1-9cee4"
BASE_URL = "https://ark.cn-beijing.volces.com/api/v3"
MODEL_NAME = "ep-m-20260517183251-qv9zq"
# ========================================================================

client = OpenAI(api_key=API_KEY, base_url=BASE_URL)

# 12 个办公功能（完整提示词）
FUNCTIONS = {
    "1. 职场日报/周报/月报自动生成":
        "你是资深职场文员，根据用户提供工作内容，自动生成标准规范日报、周报、月报，分工作完成、现存问题、后续工作计划，排版工整适合直接上交",

    "2. 企业内部通知公告撰写":
        "专业行政文案，撰写正规企业内部通知、公告、放假通知、人事通知、活动通知，格式标准、语言正式简洁、要素齐全",

    "3. 商务正式邮件自动撰写":
        "精通职场商务邮件，撰写对外对接、合作、回访、请假、汇报类正式邮件，格式标准、语气得体专业，可直接复制发送",

    "4. 会议录音文字转结构化纪要":
        "用户粘贴会议口语化录音文字，帮其梳理提炼成标准结构化会议纪要，区分参会人、会议主题、讨论内容、内容重点、达成决议、待办事项、责任人与时间",

    "5. 杂乱Excel数据清洗":
        "用户给出Excel表格，帮忙筛查问题并给出完整的处理结果版本",

    "6. PDF长文档提炼核心要点":
        "用户粘贴PDF全文/段落内容，快速精简提炼全文核心主旨、重点信息、关键结论、核心数据，去除冗余内容",

    "7. 小红书/短视频爆款文案生成":
        "擅长新媒体流量文案，根据产品、主题、场景生成吸睛标题+正文，风格贴合小红书、抖音短视频，情绪到位易出圈",

    "8. 批量数据分析+透视表制作指导":
        "针对批量业务数据，给出数据分析思路、统计维度、数据分类方法，附带Excel数据透视表详细制作流程与字段搭配技巧",

    "9. 文档内容一键转PPT文案":
        "把用户长篇文字内容，拆分梳理成标准PPT每页大纲+每页精简文案，支持简约商务/正式汇报风格，直接套用做PPT",

    "10. 表格数据分析+数据结论输出":
        "根据用户表格内数据内容，进行理性数据分析，总结数据趋势、好坏表现、存在问题、优化建议，输出完整专业数据结论",

    "11. 专业短视频/宣传片脚本撰写":
        "按照标准影视脚本格式撰写，包含镜头画面、台词文案、时长安排、背景音乐、场景动作，结构完整可直接拍摄使用",

    "12. 公众号爆款推文文章生成":
        "根据用户需求撰写公众号优质原创推文，逻辑流畅、段落分明、标题吸睛、排版舒服，符合大众阅读习惯，自带传播感"
}


# ---------------------- 文件读取功能 ----------------------
def read_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


def read_excel(file):
    df = pd.read_excel(file)
    return df.to_string(), df


def read_word(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])


# ---------------------- 文件生成功能 ----------------------
def generate_excel(content):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df = pd.DataFrame({"AI生成结果": [content]})
        df.to_excel(writer, index=False, sheet_name="结果")
    output.seek(0)
    return output


def generate_word(content):
    doc = Document()
    doc.add_heading("AI办公生成结果", 0)
    doc.add_paragraph(content)
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def generate_ppt(content):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "AI办公生成结果"
    tf = body.text_frame
    tf.text = content
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(14)
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ---------------------- AI生成 ----------------------
def generate_content(selected_func, user_input, file_content=""):
    if not user_input and not file_content:
        return "⚠️ 请输入内容或上传文件！", None, None, None

    prompt = FUNCTIONS[selected_func]
    full_input = f"文件内容：\n{file_content}\n\n用户需求：{user_input}" if file_content else user_input

    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[{"role": "system", "content": prompt},
                      {"role": "user", "content": full_input}],
            temperature=0.6, max_tokens=4000
        )
        result = response.choices[0].message.content

        excel_file = generate_excel(result)
        word_file = generate_word(result)
        ppt_file = generate_ppt(result)

        return result, excel_file, word_file, ppt_file
    except Exception as e:
        return f"❌ 错误：{str(e)}", None, None, None


# ---------------------- 界面 ----------------------
st.set_page_config(page_title="全能办公AI工具箱", layout="wide")
st.title("🧰 12合1全能办公AI工具｜文件上传 + 多格式导出")
st.markdown("### ✅ 支持 PDF / Excel / Word 上传 | 导出 Excel + Word + PPT")

# 功能选择
selected = st.selectbox("选择功能", list(FUNCTIONS.keys()))

# 文件上传
upload_file = st.file_uploader("上传文件（可选）", type=["pdf", "xlsx", "xls", "docx"])
file_content = ""

if upload_file:
    try:
        if upload_file.name.endswith(".pdf"):
            file_content = read_pdf(upload_file)
            st.success("✅ PDF 读取成功")
        elif upload_file.name.endswith((".xlsx", ".xls")):
            file_content, _ = read_excel(upload_file)
            st.success("✅ Excel 读取成功")
        elif upload_file.name.endswith(".docx"):
            file_content = read_word(upload_file)
            st.success("✅ Word 读取成功")
    except:
        st.error("❌ 文件读取失败")

# 输入需求
user_input = st.text_area("输入你的需求（必选）", height=200)

# 生成
if st.button("🚀 一键生成 + 导出全部文件", type="primary"):
    with st.spinner("AI 处理中..."):
        result, excel_file, word_file, ppt_file = generate_content(selected, user_input, file_content)
        st.success("✅ 生成完成！")
        st.markdown("---")
        st.markdown("## 📝 结果预览")
        st.write(result)

        col1, col2, col3 = st.columns(3)
        with col1:
            if excel_file:
                st.download_button("📥 下载 Excel", excel_file, "AI结果.xlsx", "application/vnd.ms-excel")
        with col2:
            if word_file:
                st.download_button("📥 下载 Word", word_file, "AI结果.docx",
                                   "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        with col3:
            if ppt_file:
                st.download_button("📥 下载 PPT", ppt_file, "AI结果.pptx",
                                   "application/vnd.openxmlformats-officedocument.presentationml.presentation")